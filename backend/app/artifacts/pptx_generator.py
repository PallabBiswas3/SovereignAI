from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


class PptxGenerator:
    def generate(self, output: Path, title: str, slides: list[dict[str, object]]) -> Path:
        output.parent.mkdir(parents=True, exist_ok=True)
        presentation = Presentation()
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = title
        title_slide.placeholders[1].text = "Generated locally by SovereignAI"
        for item in slides:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = str(item.get("title", "Briefing"))
            frame = slide.placeholders[1].text_frame
            frame.clear()
            for index, bullet in enumerate(item.get("bullets", [])):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = str(bullet)
                paragraph.level = 0
        presentation.save(output)
        return output

